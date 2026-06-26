from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colors
DEEP_BLUE   = RGBColor(0x1A, 0x23, 0x5E)   # slide backgrounds / headers
SAFFRON     = RGBColor(0xFF, 0x99, 0x00)   # accent / highlight
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF2, 0xF4, 0xF8)
MED_GRAY    = RGBColor(0x66, 0x6E, 0x8A)
BOX_BLUE    = RGBColor(0x1E, 0x3A, 0x8A)   # dark-blue boxes
BOX_ORANGE  = RGBColor(0xEA, 0x58, 0x0C)   # orange step box
BOX_GREEN   = RGBColor(0x06, 0x6B, 0x3C)   # green step box
BOX_PURPLE  = RGBColor(0x4C, 0x1D, 0x95)   # purple step box

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def bg(slide, color=DEEP_BLUE):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color, border_color=None, border_pt=0):
    """Add a filled rectangle."""
    from pptx.util import Pt as Pt2
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = Pt2(border_pt)
    else:
        shp.line.fill.background()
    return shp


def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def arrow(slide, l, t, w=0.5, h=0.3, color=SAFFRON):
    """Right-pointing arrow (chevron shape id=13)."""
    shp = slide.shapes.add_shape(13, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def step_box(slide, number, label, sublabel, l, t, w=2.2, h=1.5, box_color=BOX_BLUE):
    r = rect(slide, l, t, w, h, box_color)
    txt(slide, number, l + 0.08, t + 0.08, 0.5, 0.45, size=22, bold=True, color=SAFFRON)
    txt(slide, label, l + 0.08, t + 0.5, w - 0.16, 0.45, size=13, bold=True, color=WHITE)
    txt(slide, sublabel, l + 0.08, t + 0.95, w - 0.16, 0.48, size=10, color=RGBColor(0xCC, 0xD6, 0xFF))


def code_box(slide, code_lines, l, t, w, h):
    r = rect(slide, l, t, w, h, RGBColor(0x0D, 0x1B, 0x2A))
    box = slide.shapes.add_textbox(Inches(l + 0.15), Inches(t + 0.15),
                                   Inches(w - 0.3), Inches(h - 0.3))
    box.word_wrap = False
    tf = box.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9.5)
        run.font.color.rgb = RGBColor(0xA8, 0xD8, 0xEA)
        run.font.name = "Courier New"


# ── SLIDE 1 – Title ───────────────────────────────────────────────────────────

s1 = prs.slides.add_slide(BLANK)
bg(s1)

# decorative accent bar
rect(s1, 0, 0, 13.33, 0.12, SAFFRON)
rect(s1, 0, 7.38, 13.33, 0.12, SAFFRON)

# side accent
rect(s1, 0, 0.12, 0.18, 7.26, BOX_BLUE)

txt(s1, "RAG-Powered Q&A", 0.5, 1.6, 12.0, 1.3,
    size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s1, "How Your App Answers Questions from a Live MySQL Database",
    0.5, 3.0, 12.0, 0.8, size=22, color=SAFFRON, align=PP_ALIGN.CENTER)
txt(s1, "Retrieval-Augmented Generation  ·  Claude Opus  ·  SQLAlchemy",
    0.5, 4.0, 12.0, 0.6, size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)

# small tag bottom-right
txt(s1, "krishna_life database", 10.5, 6.9, 2.6, 0.4,
    size=10, color=MED_GRAY, align=PP_ALIGN.RIGHT)


# ── SLIDE 2 – What is RAG? ────────────────────────────────────────────────────

s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, 13.33, 0.12, SAFFRON)

txt(s2, "What is RAG?", 0.4, 0.25, 12.0, 0.7,
    size=32, bold=True, color=WHITE)

# Three concept cards
cards = [
    ("Retrieval", "Fetch real data from an\nexternal knowledge source\n(your MySQL database)",   BOX_BLUE),
    ("Augmented", "Inject that retrieved data\ninto the LLM's prompt so\nit has accurate context", BOX_ORANGE),
    ("Generation", "LLM generates a natural-\nlanguage answer grounded\nin the retrieved facts",  BOX_GREEN),
]

for i, (title, body, color) in enumerate(cards):
    lx = 0.6 + i * 4.2
    rect(s2, lx, 1.2, 3.9, 3.0, color)
    txt(s2, title, lx + 0.2, 1.35, 3.5, 0.55, size=22, bold=True, color=SAFFRON)
    txt(s2, body,  lx + 0.2, 2.0,  3.5, 1.8,  size=14, color=WHITE)
    # large letter
    txt(s2, title[0], lx + 2.8, 1.28, 0.9, 0.7, size=38, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF) if color != BOX_ORANGE else RGBColor(0xFF, 0xCC, 0x80))

# Arrow connectors between cards
for i in range(2):
    lx = 0.6 + i * 4.2 + 3.9
    txt(s2, "→", lx + 0.05, 2.45, 0.3, 0.5, size=28, bold=True, color=SAFFRON, align=PP_ALIGN.CENTER)

txt(s2,
    "Traditional LLMs only know what they were trained on.\n"
    "RAG lets them answer questions about YOUR live data — no fine-tuning needed.",
    0.6, 4.5, 12.1, 0.9, size=15, color=LIGHT_GRAY)


# ── SLIDE 3 – Architecture Overview ──────────────────────────────────────────

s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, 13.33, 0.12, SAFFRON)

txt(s3, "Architecture Overview", 0.4, 0.25, 12.0, 0.65, size=32, bold=True, color=WHITE)

# Component boxes  [label, sub, l, t, w, h, color]
components = [
    ("User",          "CLI input",          0.3,  2.8, 1.7, 1.3, BOX_PURPLE),
    ("app.py",        "Orchestrator",       2.4,  2.8, 2.0, 1.3, BOX_BLUE),
    ("Claude\nOpus",  "SQL generator\n& answerer", 5.2, 2.8, 2.2, 1.3, BOX_ORANGE),
    ("MySQL\nDB",     "krishna_life",       8.2,  2.8, 2.0, 1.3, BOX_GREEN),
    ("Schema\nCache", "connections.yml",   10.6,  2.8, 2.0, 1.3, RGBColor(0x4A, 0x55, 0x68)),
]

for label, sub, l, t, w, h, color in components:
    rect(s3, l, t, w, h, color)
    txt(s3, label, l + 0.1, t + 0.1, w - 0.2, 0.7, size=15, bold=True, color=WHITE)
    txt(s3, sub,   l + 0.1, t + 0.75, w - 0.2, 0.45, size=10, color=RGBColor(0xCC, 0xD6, 0xFF))

# Arrows (drawn as text "→")
arrows = [
    (2.05, 3.3, "question"),
    (4.45, 3.3, "schema +\nquestion"),
    (7.45, 3.3, "SQL query"),
    (10.25, 3.3, "schema\nload"),
]
for (lx, ty, label) in arrows:
    txt(s3, "→", lx, ty, 0.35, 0.5, size=20, bold=True, color=SAFFRON, align=PP_ALIGN.CENTER)
    txt(s3, label, lx - 0.1, ty + 0.45, 0.6, 0.55, size=8, color=MED_GRAY, align=PP_ALIGN.CENTER)

# Return arrows
txt(s3, "SQL\nresult ↓", 6.1, 4.3, 1.2, 0.55, size=9, color=SAFFRON, align=PP_ALIGN.CENTER)
txt(s3, "rows ↑", 9.0, 4.3, 0.9, 0.4, size=9, color=SAFFRON, align=PP_ALIGN.CENTER)

# conversation history note
rect(s3, 0.3, 5.1, 3.5, 0.9, RGBColor(0x1E, 0x1E, 0x3F))
txt(s3, "Conversation History (multi-turn)\nPrior Q&A pairs kept in memory for context",
    0.45, 5.15, 3.2, 0.75, size=10, color=RGBColor(0xCC, 0xD6, 0xFF))

txt(s3, "Prompt Cache", 5.2, 5.1, 2.2, 0.35, size=10, bold=True, color=SAFFRON)
txt(s3, "Schema sent once with cache_control: ephemeral\n→ saves tokens & latency on repeated calls",
    5.2, 5.45, 5.5, 0.6, size=10, color=LIGHT_GRAY)


# ── SLIDE 4 – Step-by-Step Flow ───────────────────────────────────────────────

s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, 13.33, 0.12, SAFFRON)

txt(s4, "Step-by-Step: How a Question Gets Answered", 0.4, 0.25, 12.5, 0.65,
    size=28, bold=True, color=WHITE)

steps = [
    ("1", "User Asks",       'Type a natural-language question\ne.g. "How many donors gave\nlast month?"',        0.25,  1.2,  BOX_PURPLE),
    ("2", "Schema Loaded",   "Database schema read once at\nstartup via SHOW CREATE TABLE.\nCached with Claude.",  2.65,  1.2,  BOX_BLUE),
    ("3", "SQL Generated",   "Claude Opus receives schema +\nquestion, returns a valid\nMySQL SELECT statement.",  5.05,  1.2,  BOX_ORANGE),
    ("4", "Query Executed",  "SQLAlchemy runs the SQL against\nthe live krishna_life database\nand returns rows.",  7.45,  1.2,  BOX_GREEN),
    ("5", "Answer Formed",   "Claude receives question + rows,\nwrites a plain-English answer.\nHistory updated.",  9.85,  1.2,  BOX_PURPLE),
]

for num, label, desc, lx, ty, color in steps:
    step_box(s4, num, label, desc, lx, ty, w=2.25, h=2.1, box_color=color)

# arrows between steps
for i in range(4):
    lx = 0.25 + (i + 1) * 2.4 - 0.15
    txt(s4, "→", lx, 2.0, 0.3, 0.5, size=20, bold=True, color=SAFFRON, align=PP_ALIGN.CENTER)

# example question / SQL / answer showcase
rect(s4, 0.25, 3.6, 12.5, 3.6, RGBColor(0x0F, 0x1C, 0x3F))
txt(s4, "Example walkthrough", 0.45, 3.7, 4.0, 0.4, size=12, bold=True, color=SAFFRON)

code_box(s4, [
    'You:  "How many active members joined in 2024?"',
    "",
    "SQL:  SELECT COUNT(*) AS total",
    "      FROM members",
    "      WHERE status = 'active'",
    "      AND YEAR(created_at) = 2024;",
], 0.35, 4.1, 5.8, 2.8)

rect(s4, 6.4, 4.1, 0.25, 2.8, RGBColor(0xFF, 0x99, 0x00))   # orange divider

txt(s4, "DB Result", 6.75, 4.1, 2.0, 0.35, size=11, bold=True, color=SAFFRON)
code_box(s4, ["[{'total': 342}]"], 6.75, 4.5, 2.5, 0.7)

txt(s4, "Claude's Answer", 9.45, 4.1, 3.5, 0.35, size=11, bold=True, color=SAFFRON)
txt(s4,
    '"342 active members joined your\ncommunity in 2024."',
    9.45, 4.5, 3.3, 1.2, size=12, color=WHITE)


# ── SLIDE 5 – The Two Claude Calls ───────────────────────────────────────────

s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, 13.33, 0.12, SAFFRON)

txt(s5, "Two Claude API Calls per Question", 0.4, 0.25, 12.5, 0.65,
    size=32, bold=True, color=WHITE)

# Call 1 box
rect(s5, 0.3, 1.1, 5.9, 5.8, BOX_BLUE)
txt(s5, "Call 1  — SQL Generation", 0.5, 1.2, 5.5, 0.5, size=17, bold=True, color=SAFFRON)
code_box(s5, [
    'system: "You are a MySQL expert.',
    "         Given schema and a question,",
    '         return ONLY valid SQL."',
    "         cache_control: ephemeral   # ← cached",
    "",
    'user:   "<natural language question>"',
    "",
    "→ Claude returns raw SQL string",
], 0.4, 1.75, 5.7, 3.2)
txt(s5, "Why cache the schema?\nThe schema is the same for every question.\nCaching it with cache_control: ephemeral\nsaves input tokens and reduces latency on\nevery subsequent call.",
    0.45, 4.95, 5.6, 1.8, size=11, color=LIGHT_GRAY)

# Call 2 box
rect(s5, 6.8, 1.1, 6.2, 5.8, BOX_GREEN)
txt(s5, "Call 2  — Plain-English Answer", 7.0, 1.2, 5.8, 0.5, size=17, bold=True, color=SAFFRON)
code_box(s5, [
    'system: "You are a data analyst.',
    "         Given a question and its SQL",
    '         result, answer concisely."',
    "",
    'user:   "Question: <question>',
    "        SQL result: <rows from DB>",
    '        <prior conversation history>"',
    "",
    "→ Claude returns natural-language answer",
], 6.9, 1.75, 6.0, 3.2)
txt(s5, "Conversation history is appended so the user can\nask follow-up questions referencing earlier answers.",
    6.9, 4.95, 6.0, 1.0, size=11, color=LIGHT_GRAY)

# divider
rect(s5, 6.4, 1.1, 0.08, 5.8, SAFFRON)


# ── SLIDE 6 – Multi-DB & Config ───────────────────────────────────────────────

s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, 13.33, 0.12, SAFFRON)

txt(s6, "Multi-Database Support & Configuration", 0.4, 0.25, 12.5, 0.65,
    size=30, bold=True, color=WHITE)

# connections.yml display
txt(s6, "connections.yml", 0.4, 1.1, 4.0, 0.45, size=16, bold=True, color=SAFFRON)
code_box(s6, [
    "default: krishna_life",
    "",
    "connections:",
    "  krishna_life:",
    "    url: mysql+pymysql://root@localhost/local_krishna_life",
    "  # donation_app:",
    "  #   url: mysql+pymysql://user:pass@host/donation_db",
], 0.4, 1.6, 6.0, 2.8)

txt(s6, "How it works", 7.0, 1.1, 5.9, 0.45, size=16, bold=True, color=SAFFRON)
features = [
    ("Connection pooling", "SQLAlchemy engines cached in _engines{} dict — one engine per URL"),
    ("Schema per-DB",      "get_schema() introspects each database independently via SHOW CREATE TABLE"),
    ("Default connection", "resolve_connection(None) returns the 'default' key from YAML config"),
    ("Lazy loading",       "_load_config() reads YAML once and caches it in _config global"),
]
for i, (title, desc) in enumerate(features):
    ty = 1.6 + i * 1.1
    rect(s6, 7.0, ty, 6.0, 0.95, RGBColor(0x1E, 0x3A, 0x5F))
    txt(s6, title, 7.15, ty + 0.06, 5.7, 0.38, size=12, bold=True, color=SAFFRON)
    txt(s6, desc,  7.15, ty + 0.45, 5.7, 0.42, size=10, color=LIGHT_GRAY)

txt(s6,
    "To add a new database:\n"
    "1. Add a new entry under connections: in connections.yml\n"
    "2. Optionally set it as default, or pass the name explicitly to resolve_connection()",
    0.4, 4.65, 12.0, 1.5, size=13, color=LIGHT_GRAY)


# ── SLIDE 7 – Key Benefits & Summary ─────────────────────────────────────────

s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, 13.33, 0.12, SAFFRON)

txt(s7, "Key Benefits & Summary", 0.4, 0.25, 12.5, 0.65, size=32, bold=True, color=WHITE)

benefits = [
    ("No fine-tuning",     "Works with any MySQL database out of the box.\nJust point it at your schema — no model training required."),
    ("Always up-to-date",  "Queries run live against the database.\nAnswers reflect the current state of your data."),
    ("Multi-turn memory",  "Conversation history lets users ask follow-up questions\nnaturally, like talking to a data analyst."),
    ("Token efficiency",   "Schema is sent with cache_control: ephemeral so it's cached\nbetween calls — saving cost and reducing latency."),
    ("Multi-DB ready",     "connections.yml lets you add databases without\nchanging any application code."),
    ("Secure & scoped",    "Claude only generates SELECT queries based on the schema.\nNo direct user input is passed to SQL execution."),
]

cols = 2
for i, (title, desc) in enumerate(benefits):
    col = i % cols
    row = i // cols
    lx = 0.4 + col * 6.4
    ty = 1.2 + row * 1.75
    rect(s7, lx, ty, 6.0, 1.55, RGBColor(0x1A, 0x33, 0x6B))
    txt(s7, f"✓  {title}", lx + 0.18, ty + 0.1, 5.6, 0.45, size=14, bold=True, color=SAFFRON)
    txt(s7, desc,           lx + 0.18, ty + 0.55, 5.6, 0.85, size=11, color=LIGHT_GRAY)

# Footer
rect(s7, 0, 6.9, 13.33, 0.6, RGBColor(0x0D, 0x18, 0x44))
txt(s7,
    "User Question  →  Schema Retrieval  →  SQL Generation  →  DB Execution  →  Plain-English Answer",
    0.3, 6.95, 12.7, 0.45, size=12, bold=True, color=SAFFRON, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────

out = "RAG_QA_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
